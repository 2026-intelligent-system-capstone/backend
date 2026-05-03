from collections.abc import Iterable
from io import BytesIO
from pathlib import PurePosixPath
from zipfile import BadZipFile, ZipFile, ZipInfo

from defusedxml import ElementTree
from docx import Document
from pptx import Presentation
from pypdf import PdfReader

from app.classroom.domain.exception import (
    ClassroomMaterialIngestDomainException,
)
from app.classroom.domain.service import ClassroomMaterialExtractedChunk

MAX_CHUNK_LENGTH = 1000
CHUNK_OVERLAP = 200
MAX_EXTRACTED_CHUNK_COUNT = 2000
MAX_EXTRACTED_TEXT_CHARS = 200_000
MAX_ZIP_FILE_COUNT = 200
MAX_ZIP_TOTAL_UNCOMPRESSED_SIZE = 20 * 1024 * 1024
MAX_ZIP_MEMBER_UNCOMPRESSED_SIZE = 5 * 1024 * 1024
TEXT_FILE_EXTENSIONS = (".txt", ".md", ".csv", ".json", ".xml")
HWPX_TEXT_PREFIXES = ("Contents/", "Preview/")


def split_text(text: str) -> Iterable[str]:
    normalized = text.strip()
    if not normalized:
        return []

    chunks: list[str] = []
    start = 0
    while start < len(normalized):
        chunk_text = normalized[start : start + MAX_CHUNK_LENGTH].strip()
        if chunk_text:
            chunks.append(chunk_text)
        if start + MAX_CHUNK_LENGTH >= len(normalized):
            break
        start += MAX_CHUNK_LENGTH - CHUNK_OVERLAP
    return chunks


def extract_pdf_chunks(
    *,
    content: bytes,
    file_name: str,
) -> list[ClassroomMaterialExtractedChunk]:
    _ = file_name
    try:
        pages = PdfReader(BytesIO(content)).pages
    except Exception as exc:
        raise ClassroomMaterialIngestDomainException(
            message=(
                "PDF 파일을 열 수 없습니다. 암호화되었거나 손상된 "
                "파일일 수 있습니다."
            )
        ) from exc

    chunks: list[ClassroomMaterialExtractedChunk] = []
    chunk_index = 0
    for page_number, page in enumerate(pages, start=1):
        text = page.extract_text()
        if not text or not text.strip():
            continue
        for chunk_text in split_text(text):
            chunks.append(
                ClassroomMaterialExtractedChunk(
                    text=chunk_text,
                    source_type="pdf",
                    source_unit_type="page",
                    citation_label=f"p.{page_number}",
                    chunk_index=chunk_index,
                    source_locator={
                        "file_name": file_name,
                        "page": page_number,
                    },
                )
            )
            chunk_index += 1
    if not chunks:
        raise ClassroomMaterialIngestDomainException(
            message=(
                "PDF에서 추출할 수 있는 텍스트가 없습니다. 스캔본이나 "
                "이미지 PDF라면 텍스트가 포함된 PDF로 다시 업로드해주세요."
            )
        )
    validate_extracted_chunk_budget(chunks)
    return chunks


def extract_docx_chunks(
    *,
    content: bytes,
    file_name: str,
) -> list[ClassroomMaterialExtractedChunk]:
    _validate_zip_container(
        content=content,
        label="DOCX",
    )
    try:
        document = Document(BytesIO(content))
    except ClassroomMaterialIngestDomainException:
        raise
    except Exception as exc:
        raise ClassroomMaterialIngestDomainException(
            message="DOCX 강의 자료를 해석하지 못했습니다."
        ) from exc

    chunks: list[ClassroomMaterialExtractedChunk] = []
    chunk_index = 0
    for paragraph_number, paragraph in enumerate(document.paragraphs, start=1):
        text = paragraph.text.strip()
        if not text:
            continue
        for chunk_text in split_text(text):
            chunks.append(
                ClassroomMaterialExtractedChunk(
                    text=chunk_text,
                    source_type="docx",
                    source_unit_type="paragraph",
                    citation_label=f"{file_name} 문단 {paragraph_number}",
                    chunk_index=chunk_index,
                    source_locator={
                        "file_name": file_name,
                        "paragraph": paragraph_number,
                    },
                )
            )
            chunk_index += 1

    for table_number, table in enumerate(document.tables, start=1):
        for row_number, row in enumerate(table.rows, start=1):
            text = "\n".join(
                cell.text.strip() for cell in row.cells if cell.text.strip()
            )
            if not text:
                continue
            for chunk_text in split_text(text):
                chunks.append(
                    ClassroomMaterialExtractedChunk(
                        text=chunk_text,
                        source_type="docx",
                        source_unit_type="table_row",
                        citation_label=(
                            f"{file_name} 표 {table_number} 행 {row_number}"
                        ),
                        chunk_index=chunk_index,
                        source_locator={
                            "file_name": file_name,
                            "table": table_number,
                            "row": row_number,
                        },
                    )
                )
                chunk_index += 1
    validate_extracted_chunk_budget(chunks)
    return chunks


def extract_pptx_chunks(
    *,
    content: bytes,
    file_name: str,
) -> list[ClassroomMaterialExtractedChunk]:
    try:
        _validate_zip_container(
            content=content,
            label="PPTX",
            raise_bad_zip=True,
        )
    except BadZipFile as exc:
        raise ClassroomMaterialIngestDomainException(
            message=(
                "PPTX 파일 구조를 읽을 수 없습니다. 손상되었거나 "
                "올바른 PPTX 파일이 아닐 수 있습니다."
            )
        ) from exc
    try:
        presentation = Presentation(BytesIO(content))
    except ClassroomMaterialIngestDomainException:
        raise
    except Exception as exc:
        raise ClassroomMaterialIngestDomainException(
            message=(
                "PPTX 슬라이드를 해석하지 못했습니다. 파일을 다시 저장한 뒤 "
                "업로드해주세요."
            )
        ) from exc

    chunks: list[ClassroomMaterialExtractedChunk] = []
    chunk_index = 0
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_text_parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text_frame.text.strip()
                if text:
                    slide_text_parts.append(text)
            if getattr(shape, "has_table", False):
                table = shape.table
                for row in table.rows:
                    row_text = "\n".join(
                        cell.text.strip()
                        for cell in row.cells
                        if cell.text.strip()
                    )
                    if row_text:
                        slide_text_parts.append(row_text)
        slide_text = "\n".join(slide_text_parts).strip()
        if not slide_text:
            continue
        for chunk_text in split_text(slide_text):
            chunks.append(
                ClassroomMaterialExtractedChunk(
                    text=chunk_text,
                    source_type="pptx",
                    source_unit_type="slide",
                    citation_label=f"slide {slide_number}",
                    chunk_index=chunk_index,
                    source_locator={
                        "file_name": file_name,
                        "slide": slide_number,
                    },
                )
            )
            chunk_index += 1
    if not chunks:
        raise ClassroomMaterialIngestDomainException(
            message=(
                "PPTX에서 추출할 수 있는 텍스트가 없습니다. "
                "이미지로만 구성된 슬라이드라면 텍스트가 포함된 파일로 "
                "다시 업로드해주세요."
            )
        )
    validate_extracted_chunk_budget(chunks)
    return chunks


def extract_hwpx_chunks(
    *,
    content: bytes,
    file_name: str,
) -> list[ClassroomMaterialExtractedChunk]:
    try:
        archive = ZipFile(BytesIO(content))
    except BadZipFile as exc:
        raise ClassroomMaterialIngestDomainException(
            message="HWPX 강의 자료를 해석하지 못했습니다."
        ) from exc

    chunks: list[ClassroomMaterialExtractedChunk] = []
    chunk_index = 0
    with archive:
        _validate_zip_limits(archive, label="HWPX")
        for info in archive.infolist():
            if info.is_dir():
                continue
            normalized_name = _normalize_archive_path(
                info.filename,
                label="HWPX",
            )
            if not _is_hwpx_text_xml(normalized_name):
                continue
            try:
                raw = archive.read(info)
                text = _extract_xml_text(raw)
            except Exception as exc:
                raise ClassroomMaterialIngestDomainException(
                    message="HWPX 강의 자료를 해석하지 못했습니다."
                ) from exc
            for chunk_text in split_text(text):
                chunks.append(
                    ClassroomMaterialExtractedChunk(
                        text=chunk_text,
                        source_type="hwpx",
                        source_unit_type="xml_file",
                        citation_label=normalized_name,
                        chunk_index=chunk_index,
                        source_locator={
                            "file_name": file_name,
                            "archive_path": normalized_name,
                        },
                    )
                )
                chunk_index += 1
    validate_extracted_chunk_budget(chunks)
    return chunks


def extract_zip_chunks(
    *,
    content: bytes,
    file_name: str,
) -> list[ClassroomMaterialExtractedChunk]:
    _ = file_name
    try:
        archive = ZipFile(BytesIO(content))
    except BadZipFile as exc:
        raise ClassroomMaterialIngestDomainException(
            message="ZIP 강의 자료를 해석하지 못했습니다."
        ) from exc

    chunks: list[ClassroomMaterialExtractedChunk] = []
    chunk_index = 0
    with archive:
        _validate_zip_limits(archive)
        for info in archive.infolist():
            if info.is_dir():
                continue
            normalized_name = _normalize_archive_path(
                info.filename,
                label="ZIP",
            )
            if not _is_text_file_name(normalized_name):
                continue
            try:
                raw = archive.read(info)
            except Exception as exc:
                raise ClassroomMaterialIngestDomainException(
                    message="ZIP 강의 자료를 해석하지 못했습니다."
                ) from exc
            text = raw.decode("utf-8", errors="ignore").strip()
            for chunk_text in split_text(text):
                chunks.append(
                    ClassroomMaterialExtractedChunk(
                        text=chunk_text,
                        source_type="zip_text",
                        source_unit_type="file",
                        citation_label=normalized_name,
                        chunk_index=chunk_index,
                        source_locator={"archive_path": normalized_name},
                    )
                )
                chunk_index += 1
    validate_extracted_chunk_budget(chunks)
    return chunks


def _validate_zip_limits(archive: ZipFile, *, label: str = "ZIP") -> None:
    _validate_zip_infos(
        infos=[info for info in archive.infolist() if not info.is_dir()],
        label=label,
    )


def _validate_zip_container(
    *,
    content: bytes,
    label: str,
    raise_bad_zip: bool = False,
) -> None:
    try:
        archive = ZipFile(BytesIO(content))
    except BadZipFile as exc:
        if raise_bad_zip:
            raise
        raise ClassroomMaterialIngestDomainException(
            message=f"{label} 강의 자료를 해석하지 못했습니다."
        ) from exc

    with archive:
        infos = [info for info in archive.infolist() if not info.is_dir()]
        _validate_zip_infos(infos=infos, label=label)


def _validate_zip_infos(*, infos: list[ZipInfo], label: str) -> None:
    if len(infos) > MAX_ZIP_FILE_COUNT:
        raise ClassroomMaterialIngestDomainException(
            message=f"{label} 강의 자료의 파일 수가 허용 범위를 초과했습니다."
        )

    total_uncompressed_size = 0
    for info in infos:
        _normalize_archive_path(info.filename, label=label)
        if info.file_size > MAX_ZIP_MEMBER_UNCOMPRESSED_SIZE:
            raise ClassroomMaterialIngestDomainException(
                message=(
                    f"{label} 강의 자료의 개별 파일 크기가 허용 범위를 "
                    "초과했습니다."
                )
            )
        total_uncompressed_size += info.file_size

    if total_uncompressed_size > MAX_ZIP_TOTAL_UNCOMPRESSED_SIZE:
        raise ClassroomMaterialIngestDomainException(
            message=(
                f"{label} 강의 자료의 전체 압축 해제 크기가 허용 범위를 "
                "초과했습니다."
            )
        )


def _normalize_archive_path(raw_path: str, *, label: str) -> str:
    normalized = raw_path.replace("\\", "/")
    path = PurePosixPath(normalized)
    if (
        path.is_absolute()
        or any(part in {"", ".", ".."} for part in path.parts)
        or any(ord(character) < 32 for character in normalized)
    ):
        raise ClassroomMaterialIngestDomainException(
            message=f"{label} 강의 자료 경로가 올바르지 않습니다."
        )
    return path.as_posix()


def validate_extracted_chunk_budget(
    chunks: list[ClassroomMaterialExtractedChunk],
) -> None:
    if len(chunks) > MAX_EXTRACTED_CHUNK_COUNT:
        raise ClassroomMaterialIngestDomainException(
            message="강의 자료에서 추출된 텍스트가 허용 범위를 초과했습니다."
        )
    total_chars = sum(len(chunk.text) for chunk in chunks)
    if total_chars > MAX_EXTRACTED_TEXT_CHARS:
        raise ClassroomMaterialIngestDomainException(
            message="강의 자료에서 추출된 텍스트가 허용 범위를 초과했습니다."
        )


def _is_text_file_name(file_name: str) -> bool:
    return file_name.lower().endswith(TEXT_FILE_EXTENSIONS)


def _is_hwpx_text_xml(file_name: str) -> bool:
    return file_name.endswith(".xml") and file_name.startswith(
        HWPX_TEXT_PREFIXES
    )


def _extract_xml_text(raw: bytes) -> str:
    root = ElementTree.fromstring(raw)
    text_parts = [text.strip() for text in root.itertext() if text.strip()]
    return "\n".join(text_parts)
